from fastapi import FastAPI
from fastapi.responses import StreamingResponse, JSONResponse
import requests
import io
import os
import cohere
from uuid import uuid4
from datetime import datetime
import textwrap
from urllib.parse import urlparse

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from docx import Document
from pptx import Presentation
import pandas as pd

import psycopg2
import psycopg2.extras


# -------------------------------------------------------------------
# Config
# -------------------------------------------------------------------

COHERE_API_KEY = os.getenv("COHERE_API_KEY")
DATABASE_URL = os.getenv("DATABASE_URL")

BASE_URL = os.getenv(
    "BASE_URL",
    "https://apex-ai-api.onrender.com"
)

if not COHERE_API_KEY:
    raise RuntimeError("COHERE_API_KEY environment variable is missing")

if not DATABASE_URL:
    raise RuntimeError("DATABASE_URL environment variable is missing")

co = cohere.Client(COHERE_API_KEY)

app = FastAPI(title="APEX AI FSD Generator")


# -------------------------------------------------------------------
# Database helpers
# -------------------------------------------------------------------

def get_db_connection():
    return psycopg2.connect(DATABASE_URL)


def init_db():
    conn = get_db_connection()
    cur = conn.cursor()

    cur.execute("""
        CREATE TABLE IF NOT EXISTS fsd_documents (
            file_id TEXT PRIMARY KEY,
            project_id TEXT,
            project_name TEXT,
            file_url TEXT,
            file_type TEXT,
            fsd_output TEXT NOT NULL,
            pdf_bytes BYTEA NOT NULL,
            created_at TIMESTAMP NOT NULL
        );
    """)

    conn.commit()
    cur.close()
    conn.close()


@app.on_event("startup")
def startup_event():
    init_db()


def save_document_to_db(
    file_id,
    project_id,
    project_name,
    file_url,
    file_type,
    fsd_output,
    pdf_bytes
):
    conn = get_db_connection()
    cur = conn.cursor()

    cur.execute("""
        INSERT INTO fsd_documents (
            file_id,
            project_id,
            project_name,
            file_url,
            file_type,
            fsd_output,
            pdf_bytes,
            created_at
        )
        VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
    """, (
        file_id,
        project_id,
        project_name,
        file_url,
        file_type,
        fsd_output,
        psycopg2.Binary(pdf_bytes),
        datetime.utcnow()
    ))

    conn.commit()
    cur.close()
    conn.close()


def get_document_from_db(file_id):
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    cur.execute("""
        SELECT
            file_id,
            project_id,
            project_name,
            file_url,
            file_type,
            fsd_output,
            pdf_bytes,
            created_at
        FROM fsd_documents
        WHERE file_id = %s
    """, (file_id,))

    row = cur.fetchone()

    cur.close()
    conn.close()

    return row


def update_pdf_in_db(file_id, pdf_bytes):
    conn = get_db_connection()
    cur = conn.cursor()

    cur.execute("""
        UPDATE fsd_documents
        SET pdf_bytes = %s
        WHERE file_id = %s
    """, (
        psycopg2.Binary(pdf_bytes),
        file_id
    ))

    conn.commit()
    cur.close()
    conn.close()


# -------------------------------------------------------------------
# File helpers
# -------------------------------------------------------------------

def get_file_extension(file_url):
    parsed = urlparse(file_url)
    path = parsed.path

    if "." not in path:
        return ""

    return path.split(".")[-1].lower()


def extract_text_from_file(file_content, file_type):
    text_content = ""

    if file_type == "docx":
        doc = Document(io.BytesIO(file_content))
        text_content = "\n".join(
            p.text for p in doc.paragraphs if p.text.strip()
        )

    elif file_type == "pptx":
        prs = Presentation(io.BytesIO(file_content))

        slide_texts = []

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_texts.append(f"\nSlide {slide_index}:")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

        text_content = "\n".join(slide_texts)

    elif file_type in ["xlsx", "xls"]:
        engine = "openpyxl" if file_type == "xlsx" else "xlrd"

        excel_data = pd.read_excel(
            io.BytesIO(file_content),
            sheet_name=None,
            engine=engine
        )

        sheet_texts = []

        for sheet_name, df in excel_data.items():
            sheet_texts.append(f"\nSheet: {sheet_name}")
            sheet_texts.append(df.to_string(index=False))

        text_content = "\n".join(sheet_texts)

    elif file_type in ["txt", "md"]:
        text_content = file_content.decode("utf-8", errors="ignore")

    else:
        raise ValueError("Unsupported file type")

    return text_content.strip()


# -------------------------------------------------------------------
# PDF creator
# -------------------------------------------------------------------

def create_pdf_buffer(text, project_id=None, project_name=None):
    buffer = io.BytesIO()

    c = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4

    x = 50
    y = height - 50
    line_height = 15
    max_chars_per_line = 95

    def new_page():
        nonlocal y
        c.showPage()
        y = height - 50
        c.setFont("Helvetica", 10)

    def draw_line(line, font_name="Helvetica", font_size=10):
        nonlocal y

        if y < 50:
            new_page()

        c.setFont(font_name, font_size)
        c.drawString(x, y, line)
        y -= line_height

    c.setFont("Helvetica-Bold", 15)
    c.drawString(x, y, "Functional Specification Document")
    y -= 30

    if project_id:
        draw_line(f"Project ID: {project_id}", "Helvetica", 10)

    if project_name:
        draw_line(f"Project Name: {project_name}", "Helvetica", 10)

    y -= 10

    for raw_line in text.splitlines():
        line = raw_line.strip()

        if not line:
            y -= line_height
            continue

        is_heading = (
            line.endswith(":")
            or line.startswith("#")
            or line[:2].isdigit()
            or line.lower().startswith((
                "section",
                "module",
                "requirement",
                "overview",
                "scope",
                "assumption",
                "dependency",
                "acceptance"
            ))
        )

        font_name = "Helvetica-Bold" if is_heading else "Helvetica"
        font_size = 11 if is_heading else 10

        wrapped_lines = textwrap.wrap(
            line,
            width=max_chars_per_line,
            replace_whitespace=False,
            drop_whitespace=True
        )

        for wrapped_line in wrapped_lines:
            draw_line(wrapped_line, font_name, font_size)

    c.save()
    buffer.seek(0)

    return buffer


# -------------------------------------------------------------------
# AI prompt
# -------------------------------------------------------------------

def build_fsd_prompt(project_id, project_name, text_content):
    return f"""
Generate a professional Functional Specification Document.

Use the following structure:

1. Document Overview
2. Project Details
3. Business Objective
4. Scope
5. Functional Requirements
6. User Roles and Responsibilities
7. Process Flow
8. Input Requirements
9. Output Requirements
10. Validation Rules
11. Error Handling
12. Assumptions
13. Dependencies
14. Acceptance Criteria
15. Conclusion

Project ID: {project_id}
Project Name: {project_name}

Source Content:
{text_content}
"""


# -------------------------------------------------------------------
# HOME
# -------------------------------------------------------------------

@app.get("/")
def home():
    return {
        "status": "API running",
        "service": "APEX AI FSD Generator"
    }


# -------------------------------------------------------------------
# GENERATE DOC
# -------------------------------------------------------------------

@app.post("/generate-doc")
async def generate_doc(data: dict):
    try:
        project_id = str(data.get("project_id") or "NA").strip()
        project_name = str(data.get("project_name") or "NA").strip()

        file_url = data.get("file_url")
        manual_content = data.get("manual_content")

        text_content = ""
        file_type = "manual"

        # ------------------------------------------------------------
        # Option 1: Manual content from APEX
        # ------------------------------------------------------------
        if manual_content and str(manual_content).strip():
            text_content = str(manual_content).strip()
            file_type = "manual"

        # ------------------------------------------------------------
        # Option 2: File URL
        # ------------------------------------------------------------
        elif file_url and str(file_url).strip():
            file_url = str(file_url).strip()
            file_type = get_file_extension(file_url)

            supported_types = ["docx", "pptx", "xlsx", "xls", "txt", "md"]

            if file_type not in supported_types:
                return JSONResponse(
                    status_code=400,
                    content={
                        "status": "ERROR",
                        "message": "Unsupported file type. Supported types: docx, pptx, xlsx, xls, txt, md"
                    }
                )

            response = requests.get(file_url, timeout=30)
            response.raise_for_status()

            file_content = response.content
            text_content = extract_text_from_file(file_content, file_type)

        else:
            return JSONResponse(
                status_code=400,
                content={
                    "status": "ERROR",
                    "message": "Either manual_content or file_url is required"
                }
            )

        # Limit AI input size
        text_content = text_content[:6000]

        if not text_content.strip():
            return JSONResponse(
                status_code=400,
                content={
                    "status": "ERROR",
                    "message": "No readable content found"
                }
            )

        # ------------------------------------------------------------
        # AI CALL - ONLY HERE
        # ------------------------------------------------------------
        prompt = build_fsd_prompt(
            project_id=project_id,
            project_name=project_name,
            text_content=text_content
        )

        ai_response = co.chat(
            model="command-a-03-2025",
            message=prompt,
            temperature=0.3
        )

        fsd_output = ai_response.text.strip()

        if not fsd_output:
            return JSONResponse(
                status_code=500,
                content={
                    "status": "ERROR",
                    "message": "AI returned empty response"
                }
            )

        # ------------------------------------------------------------
        # Create PDF once
        # ------------------------------------------------------------
        pdf_buffer = create_pdf_buffer(
            fsd_output,
            project_id,
            project_name
        )

        pdf_bytes = pdf_buffer.getvalue()

        file_id = str(uuid4())

        # ------------------------------------------------------------
        # Save permanently in PostgreSQL
        # ------------------------------------------------------------
        save_document_to_db(
            file_id=file_id,
            project_id=project_id,
            project_name=project_name,
            file_url=file_url,
            file_type=file_type,
            fsd_output=fsd_output,
            pdf_bytes=pdf_bytes
        )

        return {
            "status": "SUCCESS",
            "document": fsd_output,
            "file_id": file_id,
            "download_link": f"{BASE_URL}/download-inline/{file_id}"
        }

    except requests.exceptions.RequestException as e:
        return JSONResponse(
            status_code=400,
            content={
                "status": "ERROR",
                "message": f"File download failed: {str(e)}"
            }
        )

    except ValueError as e:
        return JSONResponse(
            status_code=400,
            content={
                "status": "ERROR",
                "message": str(e)
            }
        )

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={
                "status": "ERROR",
                "message": str(e)
            }
        )


# -------------------------------------------------------------------
# DOWNLOAD - NO AI CALL
# -------------------------------------------------------------------

@app.get("/download-inline/{file_id}")
def download_inline(file_id: str):
    try:
        metadata = get_document_from_db(file_id)

        if not metadata:
            return JSONResponse(
                status_code=404,
                content={
                    "status": "ERROR",
                    "message": "Not found"
                }
            )

        pdf_bytes = bytes(metadata["pdf_bytes"])
        pdf_buffer = io.BytesIO(pdf_bytes)
        pdf_buffer.seek(0)

        project_id = metadata.get("project_id") or "NA"
        filename = f"FSD_{project_id}.pdf"

        return StreamingResponse(
            pdf_buffer,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename}"'
            }
        )

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={
                "status": "ERROR",
                "message": str(e)
            }
        )


# -------------------------------------------------------------------
# REFRESH - NO AI CALL
# -------------------------------------------------------------------

@app.post("/refresh-doc/{file_id}")
async def refresh_doc(file_id: str):
    try:
        metadata = get_document_from_db(file_id)

        if not metadata:
            return JSONResponse(
                status_code=404,
                content={
                    "status": "ERROR",
                    "message": "Not found"
                }
            )

        pdf_buffer = create_pdf_buffer(
            metadata["fsd_output"],
            metadata["project_id"],
            metadata["project_name"]
        )

        pdf_bytes = pdf_buffer.getvalue()

        update_pdf_in_db(file_id, pdf_bytes)

        return {
            "status": "SUCCESS",
            "message": "PDF regenerated without AI call",
            "download_link": f"{BASE_URL}/download-inline/{file_id}"
        }

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={
                "status": "ERROR",
                "message": str(e)
            }
        )


# -------------------------------------------------------------------
# GET DOCUMENT TEXT - NO AI CALL
# -------------------------------------------------------------------

@app.get("/document/{file_id}")
def get_document(file_id: str):
    try:
        metadata = get_document_from_db(file_id)

        if not metadata:
            return JSONResponse(
                status_code=404,
                content={
                    "status": "ERROR",
                    "message": "Not found"
                }
            )

        return {
            "status": "SUCCESS",
            "file_id": metadata["file_id"],
            "project_id": metadata["project_id"],
            "project_name": metadata["project_name"],
            "file_url": metadata["file_url"],
            "file_type": metadata["file_type"],
            "document": metadata["fsd_output"],
            "download_link": f"{BASE_URL}/download-inline/{file_id}",
            "created_at": str(metadata["created_at"])
        }

    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={
                "status": "ERROR",
                "message": str(e)
            }
        )
